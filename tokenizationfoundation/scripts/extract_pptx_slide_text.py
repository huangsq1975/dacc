import sys
from pathlib import Path


def main() -> int:
    if len(sys.argv) < 3:
        print("Usage: python extract_pptx_slide_text.py <pptx_path> <slide_number_1_based>", file=sys.stderr)
        return 2

    pptx_path = Path(sys.argv[1])
    slide_number = int(sys.argv[2])
    if slide_number < 1:
        print("slide_number_1_based must be >= 1", file=sys.stderr)
        return 2

    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        print(f"Failed to import python-pptx. Install with: pip install python-pptx\n{e}", file=sys.stderr)
        return 1

    prs = Presentation(str(pptx_path))
    print(f"slides: {len(prs.slides)}")

    idx = slide_number - 1
    if idx >= len(prs.slides):
        print(f"No such slide: {slide_number} (max: {len(prs.slides)})", file=sys.stderr)
        return 2

    slide = prs.slides[idx]
    print(f"--- SLIDE {slide_number} TEXT ---")

    seen: set[str] = set()
    for shape in slide.shapes:
        text = getattr(shape, "text", "")
        if not isinstance(text, str):
            continue
        text = " ".join(text.split()).strip()
        if not text:
            continue
        if text in seen:
            continue
        seen.add(text)
        print(text)

    return 0


if __name__ == "__main__":
    raise SystemExit(main())

